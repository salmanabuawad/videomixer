import logging
from pathlib import Path
from typing import List

from PyPDF2 import PdfReader
from pptx import Presentation

logger = logging.getLogger(__name__)


def extract_text_from_file(path: str) -> str:
    try:
        suffix = Path(path).suffix.lower()
        if suffix == ".pdf":
            return extract_pdf(path)
        if suffix in (".pptx", ".ppt"):
            return extract_pptx(path)
        if suffix in (".txt", ".md"):
            return Path(path).read_text(encoding="utf-8", errors="ignore")
        return ""
    except Exception as e:
        logger.warning("extract_text_from_file failed for %s: %s", path, e)
        return ""


def extract_pdf(path: str) -> str:
    text_parts: List[str] = []
    try:
        reader = PdfReader(path)
    except Exception as e:
        logger.warning("PdfReader failed for %s: %s", path, e)
        return ""
    for page in reader.pages:
        try:
            text_parts.append(page.extract_text() or "")
        except Exception:
            pass
    return "\n".join(text_parts)


def extract_pptx(path: str) -> str:
    try:
        prs = Presentation(path)
    except Exception as e:
        logger.warning("Presentation failed for %s: %s", path, e)
        return ""
    text_parts: List[str] = []
    try:
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text_parts.append(shape.text)
    except Exception as e:
        logger.warning("PPTX slide walk failed for %s: %s", path, e)
    return "\n".join(text_parts)
